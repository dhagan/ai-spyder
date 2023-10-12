import docx
import pptx

# Code to extract text content from a DOCX file
def extract_text_from_docx(file_path):
    doc = docx.Document(file_path)
    text_content = ""
    for paragraph in doc.paragraphs:
        text_content += paragraph.text + "\n"
    return text_content

# Code to extract text content from a PPTX file
def extract_text_from_pptx(file_path):
    presentation = pptx.Presentation(file_path)
    text_content = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_content += shape.text + "\n"
    return text_content

# Example usage
docx_file_path = "example.docx"
pptx_file_path = "example.pptx"

docx_text = extract_text_from_docx(docx_file_path)
pptx_text = extract_text_from_pptx(pptx_file_path)

print("Text content from DOCX file:")
print(docx_text)

print("Text content from PPTX file:")
print(pptx_text)
